from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0A, 0x23, 0x40)
CREAM = RGBColor(0xF5, 0xED, 0xE1)
OLIVE = RGBColor(0x7D, 0x7F, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x15, 0x15, 0x15)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
DIM = RGBColor(0x66, 0x77, 0x88)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def txt(slide, l, t, w, h, text, font='Montserrat', sz=24, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.name = font; p.font.size = Pt(sz)
    p.font.color.rgb = clr; p.font.bold = bold; p.alignment = align
    return tb

def bar(slide, l, t, clr=CREAM):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(0.8), Pt(4))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()

def logo(slide, clr=WHITE):
    txt(slide, 0.75, 6.9, 2, 0.3, 'Navigator CFO', 'Playfair Display', 13, clr)

def snum(slide, n, clr=WHITE):
    txt(slide, 12.0, 6.9, 0.6, 0.3, f'{n:02d}', 'Montserrat', 11, clr, align=PP_ALIGN.RIGHT)

def bullets(slide, l, t, w, items, clr=WHITE, sz=18):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(len(items)*0.55))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"—  {item}"; p.font.name = 'Montserrat'; p.font.size = Pt(sz)
        p.font.color.rgb = clr; p.space_after = Pt(10)

def card(slide, l, t, w, h, fill=WHITE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()

def topbar(slide, l, t, w, clr=OLIVE):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Pt(4))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()

def img(slide, path, l, t, w):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

# ─── SLIDE 1: Title ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
bar(s, 0.75, 2.8)
txt(s, 0.75, 3.1, 8, 1.6, 'Unbreakable\nBusiness', 'Playfair Display', 60, WHITE, True)
txt(s, 0.75, 5.0, 8, 0.5, 'Financial Systems for Scalable Growth', 'Montserrat', 26, DIM)
txt(s, 0.75, 5.8, 4, 0.4, 'Presented by Lou Pacelli', 'Montserrat', 16, MUTED)
logo(s)

# ─── SLIDE 2: Objective ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, CREAM)
bar(s, 0.75, 0.8, OLIVE)
txt(s, 0.75, 1.1, 8, 0.7, 'Objective', 'Playfair Display', 44, NAVY, True)
txt(s, 0.75, 2.2, 7, 1.4, 'Design a premium business program that equips established companies to scale with financial clarity, operational discipline, and strategic control.', 'Montserrat', 20, BLACK)
s2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(4.4), Pt(4), Inches(1.0))
s2.fill.solid(); s2.fill.fore_color.rgb = OLIVE; s2.line.fill.background()
txt(s, 1.15, 4.4, 2, 0.3, 'FOCUS', 'Montserrat', 13, OLIVE, True)
txt(s, 1.15, 4.8, 7, 0.7, 'Transform finance from a reporting\nfunction into a growth engine.', 'Montserrat', 24, BLACK, True)
logo(s, NAVY); snum(s, 2, NAVY)

# ─── SLIDE 3: Who This Is For ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
bar(s, 0.75, 0.8)
txt(s, 0.75, 1.1, 5, 0.7, 'Who This Is For', 'Playfair Display', 44, WHITE, True)
bullets(s, 0.75, 2.2, 5.5, ['Established businesses ($500K–$10M revenue)', 'Owners seeking scale, not survival', 'Leaders ready for structure\nand strategic decision-making'], WHITE, 18)
txt(s, 7.5, 1.1, 5, 0.7, 'Who This Is Not For', 'Playfair Display', 32, DIM)
bullets(s, 7.5, 2.2, 5, ['Early-stage startups', 'Owners seeking theory\nwithout implementation'], DIM, 16)
logo(s); snum(s, 3)

# ─── SLIDE 4: The Problem ★ ILLUSTRATION ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, CREAM)
# Text on left half only (columns 1-6 of ~13)
bar(s, 0.75, 0.8, OLIVE)
txt(s, 0.75, 1.1, 6, 0.9, "Most Businesses Don't Fail\nfrom Lack of Revenue", 'Playfair Display', 36, NAVY, True)
txt(s, 0.75, 2.4, 6, 0.5, 'They fail from lack of financial control.', 'Montserrat', 20, OLIVE, True)
txt(s, 0.75, 3.3, 3, 0.3, 'COMMON GAPS', 'Montserrat', 13, OLIVE, True)
bullets(s, 0.75, 3.8, 6, ['No visibility into cash flow', 'Poor decisions from incomplete data', 'Growth without structure', 'Reactive instead of strategic leadership'], BLACK, 17)
# Illustration on right half
img(s, 'illustrations/slide04_financial_control.png', 8.5, 1.2, 4.0)
logo(s, NAVY); snum(s, 4, NAVY)

# ─── SLIDE 5: Framework ★ ILLUSTRATION ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
# Text on left half
bar(s, 0.75, 0.8)
txt(s, 0.75, 1.1, 7, 0.7, 'Unbreakable Business\nFramework', 'Playfair Display', 40, WHITE, True)
txt(s, 0.75, 2.4, 7, 0.5, 'A Three-Phase Model for Growth', 'Montserrat', 20, DIM)
txt(s, 0.75, 3.6, 7, 0.5, 'Foundations  →  Activation  →  Optimization', 'Montserrat', 24, CREAM, True)
txt(s, 0.75, 4.5, 7, 0.5, 'Awareness  →  Implementation  →  Strategic Mastery', 'Montserrat', 16, DIM)
# Illustration on right
img(s, 'illustrations/slide05_three_phases.png', 8.8, 1.5, 3.8)
logo(s); snum(s, 5)

# ─── SLIDE 6: Value Progression ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, CREAM)
bar(s, 0.75, 0.8, OLIVE)
txt(s, 0.75, 1.1, 8, 0.7, 'A Clear Progression of Value', 'Playfair Display', 44, NAVY, True)
for i, (ph, ti, desc) in enumerate([
    ('Phase 1','Foundations','Education and\nFinancial Clarity'),
    ('Phase 2','Activation','Systems and\nOperational Execution'),
    ('Phase 3','Optimization','Strategic Growth\nand Advisory')]):
    x = 0.75 + i * 4.0
    card(s, x, 2.5, 3.6, 3.2)
    topbar(s, x, 2.5, 3.6)
    txt(s, x+0.35, 2.9, 2.9, 0.3, ph.upper(), 'Montserrat', 11, OLIVE, True)
    txt(s, x+0.35, 3.4, 2.9, 0.5, ti, 'Montserrat', 26, BLACK, True)
    txt(s, x+0.35, 4.1, 2.9, 1.0, desc, 'Montserrat', 16, MUTED)
txt(s, 0.75, 6.0, 8, 0.3, 'Each level builds upon the last.', 'Montserrat', 14, MUTED)
logo(s, NAVY); snum(s, 6, NAVY)

# ─── SLIDE 7: Optimization Detail ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
bar(s, 0.75, 0.8)
txt(s, 0.75, 1.1, 5.5, 0.7, 'Optimization', 'Playfair Display', 44, WHITE, True)
txt(s, 0.75, 2.0, 5.5, 0.5, 'Lead at a Strategic Level', 'Montserrat', 22, CREAM, True)
bullets(s, 0.75, 2.8, 5.5, ['Strategic decision-making using financial data', 'Growth planning — hiring, expansion, pricing', 'Enterprise-level thinking'], WHITE, 18)
# Right column
txt(s, 7.5, 0.8, 2, 0.3, 'DELIVERY', 'Montserrat', 11, CREAM, True)
bullets(s, 7.5, 1.2, 5, ['One-on-one advisory', 'Select access to Theresa Pori', 'Optional in-person intensives'], WHITE, 16)
txt(s, 7.5, 3.2, 2, 0.3, 'OUTCOME', 'Montserrat', 11, CREAM, True)
txt(s, 7.5, 3.6, 5, 0.5, 'Scalable, strategically led business', 'Montserrat', 20, WHITE, True)
txt(s, 7.5, 4.5, 3, 0.7, '$22,000+', 'Playfair Display', 48, CREAM, True)
logo(s); snum(s, 7)

# ─── SLIDE 8: Value Ladder ★ ILLUSTRATION ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, CREAM)
bar(s, 0.75, 0.8, OLIVE)
txt(s, 0.75, 1.1, 6, 0.7, 'Value Ladder Strategy', 'Playfair Display', 44, NAVY, True)
txt(s, 0.75, 2.1, 6, 0.7, 'This is not three offers.\nIt is one ascension path.', 'Montserrat', 18, MUTED)
# Stacked cards on the left
for i, (label, title) in enumerate([
    ('Entry Through Clarity','Foundations'),
    ('Commitment Through Execution','Activation'),
    ('Transformation Through Proximity','Optimization')]):
    y = 3.2 + i * 1.25
    card(s, 0.75, y, 6.0, 1.0)
    txt(s, 1.1, y + 0.1, 3, 0.25, label.upper(), 'Montserrat', 10, OLIVE, True)
    txt(s, 1.1, y + 0.45, 3, 0.4, title, 'Montserrat', 22, BLACK, True)
# Illustration on the right
img(s, 'illustrations/slide08_value_ladder.png', 8.5, 1.5, 4.0)
logo(s, NAVY); snum(s, 8, NAVY)

# ─── SLIDE 9: Market Comparison ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
bar(s, 0.75, 0.8)
txt(s, 0.75, 1.1, 8, 0.7, 'Market Comparison', 'Playfair Display', 44, WHITE, True)
# Left column
txt(s, 0.75, 2.4, 5, 0.5, 'Mass Market', 'Montserrat', 26, DIM, True)
txt(s, 0.75, 3.0, 5, 0.3, 'Tony Robbins / Grant Cardone', 'Montserrat', 13, DIM)
bullets(s, 0.75, 3.5, 5, ['High volume', 'Motivational focus', 'Limited customization'], DIM, 16)
# Right column
txt(s, 7.5, 2.4, 5, 0.5, 'Unbreakable Business', 'Montserrat', 26, CREAM, True)
txt(s, 7.5, 3.0, 5, 0.3, 'Navigator CFO', 'Montserrat', 13, DIM)
bullets(s, 7.5, 3.5, 5, ['Targeted — established businesses', 'Deep financial & operational integration', 'High-touch transformation'], WHITE, 16)
txt(s, 0.75, 5.6, 10, 0.4, 'Positioning:  Premium, practical, outcome-driven', 'Montserrat', 18, CREAM, True)
logo(s); snum(s, 9)

# ─── SLIDE 10: Theresa's Role ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, CREAM)
bar(s, 0.75, 0.8, OLIVE)
txt(s, 0.75, 1.1, 5.5, 1.0, "Theresa's\nStrategic Role", 'Playfair Display', 44, NAVY, True)
txt(s, 0.75, 2.6, 5.5, 0.5, 'Leveraged, Not Diluted', 'Montserrat', 22, OLIVE, True)
# Right side
card(s, 7.5, 1.0, 5, 2.8)
bullets(s, 7.8, 1.3, 4.4, ['Present in high-impact moments', 'Strategic visibility in Activation', 'Direct involvement in Optimization'], BLACK, 18)
txt(s, 7.5, 4.4, 2, 0.3, 'EFFECT', 'Montserrat', 11, OLIVE, True)
txt(s, 7.5, 4.8, 5, 0.7, 'Increases exclusivity\nand perceived value', 'Montserrat', 22, BLACK, True)
logo(s, NAVY); snum(s, 10, NAVY)

# ─── SLIDE 11: Revenue Model ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
bar(s, 0.75, 0.8)
txt(s, 0.75, 1.1, 8, 0.7, 'Revenue Model Potential', 'Playfair Display', 44, WHITE, True)
txt(s, 0.75, 2.0, 4, 0.3, 'Illustrative Example', 'Montserrat', 13, DIM)
for i, (tier, clients, rev) in enumerate([
    ('Foundations','20','$40,000'),
    ('Activation','15','$112,500'),
    ('Optimization','8','$176,000')]):
    y = 3.0 + i * 0.75
    txt(s, 0.75, y, 3, 0.5, tier, 'Montserrat', 20, WHITE)
    txt(s, 4.5, y, 1.5, 0.5, clients, 'Montserrat', 20, DIM, align=PP_ALIGN.CENTER)
    txt(s, 6.5, y, 2.5, 0.5, rev, 'Montserrat', 20, WHITE, align=PP_ALIGN.RIGHT)
txt(s, 0.75, 5.5, 3, 0.5, 'Total per Cohort', 'Montserrat', 22, WHITE, True)
txt(s, 6.5, 5.5, 2.5, 0.5, '$328,500', 'Playfair Display', 34, CREAM, True, PP_ALIGN.RIGHT)
txt(s, 0.75, 6.3, 8, 0.3, 'Excludes upsells, renewals, and backend programs.', 'Montserrat', 13, DIM)
logo(s); snum(s, 11)

# ─── SLIDE 12: Go-To-Market ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, CREAM)
bar(s, 0.75, 0.8, OLIVE)
txt(s, 0.75, 1.1, 8, 0.7, 'Go-To-Market Strategy', 'Playfair Display', 44, NAVY, True)
for i, (phase, items) in enumerate([
    ('Phase 1 — Pilot', ['10–20 qualified business owners','Focus on Activation Tier']),
    ('Phase 2 — Refine', ['Identify objections','Strengthen outcome articulation']),
    ('Phase 3 — Scale', ['Add Foundations as feeder','Expand Optimization selectively'])]):
    x = 0.75 + i * 4.0
    card(s, x, 2.5, 3.6, 3.5)
    topbar(s, x, 2.5, 3.6)
    txt(s, x+0.35, 2.9, 2.9, 0.3, phase.upper(), 'Montserrat', 11, OLIVE, True)
    bullets(s, x+0.35, 3.5, 2.9, items, BLACK, 16)
logo(s, NAVY); snum(s, 12, NAVY)

# ─── SLIDE 13: Pricing Validation ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
bar(s, 0.75, 0.8)
txt(s, 0.75, 1.1, 5.5, 0.9, 'Pricing Validation\nStrategy', 'Playfair Display', 40, WHITE, True)
# Right side
txt(s, 7.5, 0.8, 2, 0.3, 'HOW WE TEST', 'Montserrat', 11, CREAM, True)
bullets(s, 7.5, 1.3, 5, ['Sell at full price first', 'Track conversion rates\nand objections', 'Adjust messaging before pricing'], WHITE, 17)
txt(s, 7.5, 3.8, 2, 0.3, 'KEY SIGNAL', 'Montserrat', 11, CREAM, True)
txt(s, 7.5, 4.3, 5, 0.9, 'If prospects hesitate,\nit is usually clarity —\nnot cost.', 'Montserrat', 22, CREAM, True)
logo(s); snum(s, 13)

# ─── SLIDE 14: Why This Works ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, CREAM)
bar(s, 0.75, 0.8, OLIVE)
txt(s, 0.75, 1.1, 8, 0.7, 'Why This Works', 'Playfair Display', 44, NAVY, True)
txt(s, 0.75, 2.1, 5, 0.5, 'Aligned with Market Reality', 'Montserrat', 22, OLIVE, True)
bullets(s, 0.75, 3.2, 9, [
    'Business owners pay for outcomes, not information',
    'Financial control is a universal bottleneck',
    'Structured progression increases\nretention and lifetime value'], BLACK, 20)
logo(s, NAVY); snum(s, 14, NAVY)

# ─── SLIDE 15: The Opportunity ★ ILLUSTRATION ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, NAVY)
bar(s, 0.75, 0.8)
txt(s, 0.75, 1.1, 6.5, 0.7, 'The Opportunity', 'Playfair Display', 44, WHITE, True)
txt(s, 0.75, 2.2, 6.5, 0.7, 'Unbreakable Business\ncan own this position.', 'Montserrat', 24, CREAM, True)
txt(s, 0.75, 3.6, 6.5, 0.4, 'Not just coaching. Not just education.', 'Montserrat', 18, DIM)
txt(s, 0.75, 4.4, 6.5, 0.9, 'A system for building\nbusinesses that actually scale.', 'Montserrat', 24, WHITE, True)
# Illustration right side
img(s, 'illustrations/slide15_opportunity.png', 8.5, 1.2, 4.0)
logo(s); snum(s, 15)

# ─── SLIDE 16: Closing ───
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, CREAM)
bar(s, 0.75, 3.5, OLIVE)
txt(s, 0.75, 3.8, 10, 1.2, 'This Is How Businesses\nBecome Unbreakable', 'Playfair Display', 44, NAVY, True)
txt(s, 0.75, 5.4, 3, 0.4, 'Clarity in numbers', 'Montserrat', 18, OLIVE, True)
txt(s, 4.0, 5.4, 3, 0.4, 'Disciplined systems', 'Montserrat', 18, OLIVE, True)
txt(s, 7.3, 5.4, 3, 0.4, 'Strategic leadership', 'Montserrat', 18, OLIVE, True)
txt(s, 0.75, 6.1, 8, 0.6, 'The question is not whether the market needs this.\nIt is whether we position it strongly enough to lead it.', 'Montserrat', 15, MUTED)
snum(s, 16, NAVY)

out = 'Unbreakable Business - Navigator CFO v3.pptx'
prs.save(out)
print(f'Saved: {out}')
